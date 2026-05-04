"""
services/slide_generator.py
Tạo file PPTX thật từ structured slides JSON (output của Content Service).
Dùng python-pptx để sinh PowerPoint chuyên nghiệp.
"""
import io
import logging
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# ── Colors ──
BG_CREAM = RGBColor(0xF5, 0xF0, 0xE8)
TEXT_BROWN = RGBColor(0x5D, 0x3A, 0x1A)
TEXT_LIGHT = RGBColor(0x7A, 0x5A, 0x3A)
ACCENT = RGBColor(0x8B, 0x5E, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BROWN = RGBColor(0x3E, 0x22, 0x10)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def generate_pptx(title: str, slides_data: list[dict]) -> bytes:
    """
    Tạo file PPTX từ structured slides JSON.

    Args:
        title: Tiêu đề bài thuyết trình
        slides_data: Danh sách slide objects từ Content Service

    Returns:
        bytes — nội dung file PPTX
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in sorted(slides_data, key=lambda s: s.get("slide_order", 0)):
        layout_type = slide_data.get("layout_type", "content")

        if layout_type == "title":
            _add_title_slide(prs, slide_data)
        elif layout_type == "two_column":
            _add_two_column_slide(prs, slide_data)
        elif layout_type == "timeline":
            _add_timeline_slide(prs, slide_data)
        elif layout_type == "summary":
            _add_summary_slide(prs, slide_data)
        elif layout_type == "quote":
            _add_quote_slide(prs, slide_data)
        elif layout_type == "table":
            _add_table_slide(prs, slide_data)
        elif layout_type == "section_divider":
            _add_section_divider(prs, slide_data)
        else:
            _add_content_slide(prs, slide_data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── Helpers ──

def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_BROWN, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def _add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xD5, 0xC5)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


# ── Slide Builders ──

def _add_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_bg(slide, BG_CREAM)

    # Title
    _add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(2),
                 data.get("title", ""), font_size=48, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Georgia")

    # Subtitle
    if data.get("subtitle"):
        _add_textbox(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.8),
                     data["subtitle"], font_size=18, color=TEXT_LIGHT,
                     alignment=PP_ALIGN.CENTER)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(5.5), Inches(5.2), Inches(2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Era
    if data.get("content"):
        _add_textbox(slide, Inches(3), Inches(5.5), Inches(7), Inches(0.5),
                     data["content"], font_size=14, color=ACCENT,
                     alignment=PP_ALIGN.CENTER)


def _add_content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CREAM)

    # Accent bar + heading
    _add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.55))
    _add_textbox(slide, Inches(0.85), Inches(0.45), Inches(8), Inches(0.7),
                 data.get("title", ""), font_size=28, bold=True, font_name="Georgia")

    # Content text
    y = Inches(1.4)
    if data.get("content"):
        _add_textbox(slide, Inches(0.85), y, Inches(7), Inches(0.8),
                     data["content"], font_size=13, color=TEXT_LIGHT, italic=True)
        y += Inches(0.9)

    # Bullets
    bullets = data.get("bullets", [])
    for i, bullet in enumerate(bullets[:5]):
        _add_textbox(slide, Inches(1.2), y + Inches(i * 0.55), Inches(7), Inches(0.5),
                     f"●  {bullet}", font_size=13, color=TEXT_BROWN)

    # Image placeholder on right
    card = _add_card(slide, Inches(9), Inches(1.4), Inches(3.8), Inches(5))


def _add_two_column_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CREAM)

    # Heading
    _add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.55))
    _add_textbox(slide, Inches(0.85), Inches(0.45), Inches(8), Inches(0.7),
                 data.get("title", ""), font_size=28, bold=True, font_name="Georgia")

    columns = data.get("columns", [])
    col_width = Inches(5.5)
    gap = Inches(0.5)

    for i, col in enumerate(columns[:3]):
        left = Inches(0.8) + i * (col_width + gap)
        top = Inches(1.8)

        # Card background
        _add_card(slide, left, top, col_width, Inches(4.5))

        # Icon placeholder (text)
        icon_text = col.get("icon", col.get("icon_name", "★"))
        _add_textbox(slide, left, top + Inches(0.4), col_width, Inches(0.6),
                     icon_text, font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

        # Heading
        _add_textbox(slide, left + Inches(0.3), top + Inches(1.2), col_width - Inches(0.6), Inches(0.5),
                     col.get("heading", ""), font_size=16, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Georgia")

        # Text
        _add_textbox(slide, left + Inches(0.3), top + Inches(1.8), col_width - Inches(0.6), Inches(2.2),
                     col.get("text", ""), font_size=12, color=TEXT_LIGHT,
                     alignment=PP_ALIGN.CENTER)


def _add_timeline_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CREAM)

    _add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.55))
    _add_textbox(slide, Inches(0.85), Inches(0.45), Inches(8), Inches(0.7),
                 data.get("title", ""), font_size=28, bold=True, font_name="Georgia")

    events = data.get("events", [])
    n = min(len(events), 4)
    card_w = Inches((12) / max(n, 1))

    for i, ev in enumerate(events[:4]):
        left = Inches(0.6) + i * (card_w + Inches(0.15))
        top = Inches(1.8)

        _add_card(slide, left, top, card_w - Inches(0.15), Inches(4.5))

        # Place
        place = ev.get("place", ev.get("year", ""))
        _add_textbox(slide, left + Inches(0.2), top + Inches(0.5),
                     card_w - Inches(0.5), Inches(0.5),
                     place, font_size=16, bold=True, alignment=PP_ALIGN.CENTER,
                     font_name="Georgia")

        # Year
        if ev.get("year") and ev.get("place"):
            _add_textbox(slide, left + Inches(0.2), top + Inches(1.1),
                         card_w - Inches(0.5), Inches(0.4),
                         ev["year"], font_size=11, color=ACCENT,
                         alignment=PP_ALIGN.CENTER)

        # Text
        _add_textbox(slide, left + Inches(0.2), top + Inches(1.6),
                     card_w - Inches(0.5), Inches(2.5),
                     ev.get("text", ""), font_size=11, color=TEXT_LIGHT,
                     alignment=PP_ALIGN.CENTER)


def _add_summary_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CREAM)

    # Title centered
    _add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
                 data.get("title", ""), font_size=36, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Georgia")

    if data.get("subtitle"):
        _add_textbox(slide, Inches(2), Inches(1.5), Inches(9), Inches(0.5),
                     data["subtitle"], font_size=14, color=TEXT_LIGHT,
                     italic=True, alignment=PP_ALIGN.CENTER)

    # Cards
    bullets = data.get("bullets", [])
    n = min(len(bullets), 3)
    card_w = Inches(3.5)
    total_w = n * card_w + (n - 1) * Inches(0.3)
    start_x = (SLIDE_WIDTH - total_w) / 2

    for i, b in enumerate(bullets[:3]):
        left = start_x + i * (card_w + Inches(0.3))
        top = Inches(2.5)

        _add_card(slide, left, top, card_w, Inches(3))
        _add_textbox(slide, left + Inches(0.3), top + Inches(0.5),
                     card_w - Inches(0.6), Inches(2),
                     b, font_size=13, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    if data.get("footer_text"):
        _add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                     data["footer_text"], font_size=13, color=ACCENT,
                     italic=True, alignment=PP_ALIGN.CENTER)


def _add_quote_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CREAM)

    _add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.55))
    _add_textbox(slide, Inches(0.85), Inches(0.45), Inches(8), Inches(0.7),
                 data.get("title", ""), font_size=28, bold=True, font_name="Georgia")

    # Quote
    quote = data.get("quote_text", data.get("content", ""))
    _add_textbox(slide, Inches(1.5), Inches(2), Inches(10), Inches(3),
                 f'"{quote}"', font_size=20, italic=True,
                 alignment=PP_ALIGN.CENTER, font_name="Georgia")

    # Author
    if data.get("quote_author"):
        _add_textbox(slide, Inches(2), Inches(5), Inches(9), Inches(0.5),
                     f"— {data['quote_author']}", font_size=14, color=ACCENT,
                     italic=True, alignment=PP_ALIGN.CENTER)

    if data.get("quote_context"):
        _add_textbox(slide, Inches(2), Inches(5.7), Inches(9), Inches(0.8),
                     data["quote_context"], font_size=12, color=TEXT_LIGHT,
                     alignment=PP_ALIGN.CENTER)


def _add_table_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_CREAM)

    _add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.55))
    _add_textbox(slide, Inches(0.85), Inches(0.45), Inches(8), Inches(0.7),
                 data.get("title", ""), font_size=28, bold=True, font_name="Georgia")

    headers = data.get("table_headers", [])
    rows = data.get("table_rows", [])

    if headers and rows:
        n_rows = len(rows) + 1
        n_cols = len(headers)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5 * n_rows)
        )
        table = table_shape.table

        # Header
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT

        # Rows
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = TEXT_BROWN

    if data.get("footer_text"):
        _add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                     data["footer_text"], font_size=13, color=ACCENT,
                     italic=True, alignment=PP_ALIGN.CENTER)


def _add_section_divider(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BROWN)

    # Lines
    for y in [Inches(2.8), Inches(4.8)]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(5.5), y, Inches(2), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xD4, 0xA8, 0x43)
        line.line.fill.background()

    _add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                 data.get("title", ""), font_size=36, bold=True,
                 color=BG_CREAM, alignment=PP_ALIGN.CENTER, font_name="Georgia")

    if data.get("subtitle"):
        _add_textbox(slide, Inches(2), Inches(4), Inches(9), Inches(0.6),
                     data["subtitle"], font_size=16, color=RGBColor(0xD4, 0xA8, 0x43),
                     italic=True, alignment=PP_ALIGN.CENTER)
